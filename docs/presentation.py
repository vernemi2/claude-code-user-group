#!/usr/bin/env python3
"""Generate the user group presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
ACCENT_BLUE = RGBColor(0x00, 0xA1, 0xE0)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x99)
DARK_CARD = RGBColor(0x25, 0x25, 0x3E)
CODE_BG = RGBColor(0x1E, 0x1E, 0x30)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=20,
                    color=WHITE, bullet_color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(8)
        p.space_after = Pt(4)

        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "\u25B8 "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Calibri"

        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = "Calibri"
    return txBox


def add_card(slide, left, top, width, height, color=DARK_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_card_with_title(slide, left, top, width, height, title, items,
                        accent=ACCENT_BLUE, font_size=16):
    add_card(slide, left, top, width, height)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2),
                 width - Inches(0.6), Inches(0.5), title,
                 font_size=20, color=accent, bold=True)
    add_bullet_list(slide, left + Inches(0.3), top + Inches(0.7),
                    width - Inches(0.6), height - Inches(0.9),
                    items, font_size=font_size, color=LIGHT_GRAY, bullet_color=accent)


def add_accent_bar(slide, top, color=ACCENT_ORANGE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0), top, Inches(0.15), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_phase_card(slide, left, top, width, number, title, subtitle, color):
    add_card(slide, left, top, width, Inches(1.2))
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     left + Inches(0.25), top + Inches(0.2),
                                     Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, left + Inches(1.0), top + Inches(0.15),
                 width - Inches(1.3), Inches(0.4), title,
                 font_size=20, color=WHITE, bold=True)
    add_text_box(slide, left + Inches(1.0), top + Inches(0.6),
                 width - Inches(1.3), Inches(0.5), subtitle,
                 font_size=16, color=MID_GRAY)


# ─── SLIDE 1: Title ──────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "The Developer Who Wasn't There",
             font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "Autonomous Salesforce Development with Claude Code",
             font_size=28, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

# Decorative line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(4.5), Inches(4.0), Inches(4), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             "Michal Verner",
             font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)



# ─── SLIDE 2: What if... ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1),
             "What if you could describe a feature...",
             font_size=36, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
             "...and have it built, tested, deployed, and validated",
             font_size=36, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.4), Inches(11), Inches(1),
             "without writing a single line of code?",
             font_size=36, color=WHITE, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 3: Audience Choice ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "You Decide: Pick a Feature",
             font_size=40, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "The AI agent will build your choice live \u2014 from user story to production",
             font_size=20, color=MID_GRAY)

# Card 1
add_card(slide, Inches(0.5), Inches(2.0), Inches(3.9), Inches(4.5))
add_text_box(slide, Inches(0.8), Inches(2.1), Inches(3.3), Inches(0.5),
             "A", font_size=32, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(1.3), Inches(2.15), Inches(2.8), Inches(0.5),
             "Invoice Aging", font_size=22, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(3.0), Inches(3.3), Inches(1.5),
             "As an account manager, I want to "
             "see an aging breakdown of unpaid "
             "invoices (Current, 1-30, 31-60, "
             "61-90, 90+ days overdue) on the "
             "Account page, so I can prioritize "
             "collections and assess credit risk "
             "at a glance.",
             font_size=16, color=MID_GRAY)

# Card 2
add_card(slide, Inches(4.7), Inches(2.0), Inches(3.9), Inches(4.5))
add_text_box(slide, Inches(5.0), Inches(2.1), Inches(3.3), Inches(0.5),
             "B", font_size=32, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(5.5), Inches(2.15), Inches(2.8), Inches(0.5),
             "Status Progress Tracker", font_size=22, color=WHITE, bold=True)
add_text_box(slide, Inches(5.0), Inches(3.0), Inches(3.3), Inches(1.5),
             "As an account manager, I want to see a "
             "visual progress tracker on the "
             "Invoice page showing Draft \u2192 Sent "
             "\u2192 Paid with the current step "
             "highlighted, so I can instantly "
             "understand where each invoice is "
             "in its lifecycle.",
             font_size=16, color=MID_GRAY)

# Card 3
add_card(slide, Inches(8.9), Inches(2.0), Inches(3.9), Inches(4.5))
add_text_box(slide, Inches(9.2), Inches(2.1), Inches(3.3), Inches(0.5),
             "C", font_size=32, color=ACCENT_PURPLE, bold=True)
add_text_box(slide, Inches(9.7), Inches(2.15), Inches(2.8), Inches(0.5),
             "Quick Mark as Paid", font_size=22, color=WHITE, bold=True)
add_text_box(slide, Inches(9.2), Inches(3.0), Inches(3.3), Inches(1.5),
             "As an account manager, I want a "
             "\"Mark as Paid\" button on the "
             "Invoice page that opens a confirm "
             "modal and sets Status to Paid plus "
             "today's Paid Date, so I can close "
             "out paid invoices in one click.",
             font_size=16, color=MID_GRAY)


# ─── SLIDE 4: Starting the Agent ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
             "/story-to-feature",
             font_size=52, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER,
             font_name="Consolas")
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             "One command. Zero intervention. Let's go.",
             font_size=24, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 5: While we wait... ───────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
             "While our developer works...",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
             "Let me show you how we got here.",
             font_size=24, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 6: The Foundation ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.5))
add_text_box(slide, Inches(0.5), Inches(0.4), Inches(11), Inches(0.8),
             "   The Foundation", font_size=40, color=WHITE, bold=True)
add_text_box(slide, Inches(0.5), Inches(1.2), Inches(11), Inches(0.5),
             "   What makes autonomous development possible?", font_size=20, color=MID_GRAY)

add_card_with_title(slide, Inches(0.5), Inches(2.2), Inches(5.9), Inches(4.8),
                    "Project Setup", [
                        "Salesforce DX project (API v66.0)",
                        "Developer Edition org",
                        "Pre-commit hooks (Husky + lint-staged)",
                        "ESLint + Prettier enforced",
                        "Jest for LWC unit tests",
                    ], accent=ACCENT_BLUE, font_size=18)

add_card_with_title(slide, Inches(6.8), Inches(2.2), Inches(5.9), Inches(4.8),
                    "Third-Party Libraries", [
                        "SOQL Lib \u2014 fluent query builder",
                        "DML Lib \u2014 unit of work pattern",
                        "TriggerHandler \u2014 one trigger per object",
                        "UniversalMocker \u2014 Apex Stub API",
                        "InstanceProvider \u2014 dependency injection",
                    ], accent=ACCENT_GREEN, font_size=18)


# ─── SLIDE 7: Why Libraries Matter ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.5), ACCENT_GREEN)
add_text_box(slide, Inches(0.5), Inches(0.4), Inches(11), Inches(0.8),
             "   Why Libraries Matter for AI", font_size=40, color=WHITE, bold=True)

add_card_with_title(slide, Inches(0.5), Inches(1.8), Inches(3.8), Inches(5.2),
                    "Consistency", [
                        "Every query uses the same fluent API pattern",
                        "AI learns ONE pattern, applies it everywhere",
                        "No variation = fewer bugs",
                    ], accent=ACCENT_BLUE, font_size=16)

add_card_with_title(slide, Inches(4.7), Inches(1.8), Inches(3.8), Inches(5.2),
                    "Mockability", [
                        "Every dependency can be mocked at the boundary",
                        "Zero database hits in tests",
                        "AI can write tests without org state dependencies",
                    ], accent=ACCENT_GREEN, font_size=16)

add_card_with_title(slide, Inches(8.9), Inches(1.8), Inches(3.8), Inches(5.2),
                    "Guard Rails", [
                        "Hard rules the AI can follow:",
                        "\"Never write raw SOQL\"",
                        "\"Never use raw DML\"",
                        "Libraries enforce patterns that prevent bad code",
                    ], accent=ACCENT_ORANGE, font_size=16)


# ─── SLIDE 8: CLAUDE.md ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.5), ACCENT_ORANGE)
add_text_box(slide, Inches(0.5), Inches(0.4), Inches(11), Inches(0.8),
             "   CLAUDE.md \u2014 The Instruction Manual", font_size=40, color=WHITE, bold=True)
add_text_box(slide, Inches(0.5), Inches(1.2), Inches(11), Inches(0.5),
             "   Every convention, pattern, and rule \u2014 in one file the agent always reads",
             font_size=20, color=MID_GRAY)

sections = [
    ("SF CLI Commands", "Deploy, test, open org \u2014\nready-made commands", ACCENT_BLUE),
    ("Trigger Pattern", "TriggerHandler base class,\none trigger per object", ACCENT_GREEN),
    ("Service Layer", "Instance-based, resolved\nvia InstanceProvider", ACCENT_ORANGE),
    ("SOQL / DML Libs", "Fluent API for all queries\nand DML operations", ACCENT_PURPLE),
    ("Test Conventions", "Mock everything, 95%\ncoverage, no DB hits", ACCENT_RED),
    ("Hard Rules", "No SOQL in loops, bulk-safe,\nno hardcoded IDs", RGBColor(0x00, 0xBC, 0xD4)),
]

for i, (title, desc, color) in enumerate(sections):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + Inches(4.1) * col
    top = Inches(2.0) + Inches(2.6) * row
    add_card(slide, left, top, Inches(3.8), Inches(2.3))
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2),
                 Inches(3.2), Inches(0.4), title,
                 font_size=22, color=color, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.8),
                 Inches(3.2), Inches(1.2), desc,
                 font_size=16, color=LIGHT_GRAY)


# ─── SLIDE 9: CLAUDE.md Code Example ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.5), ACCENT_ORANGE)
add_text_box(slide, Inches(0.5), Inches(0.4), Inches(11), Inches(0.8),
             "   CLAUDE.md \u2014 Teaching by Example", font_size=40, color=WHITE, bold=True)

# Left: Code example
add_card(slide, Inches(0.5), Inches(1.6), Inches(6.0), Inches(5.4), CODE_BG)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5.4), Inches(0.4),
             "Selector Pattern (from CLAUDE.md)", font_size=14, color=MID_GRAY)

code_text = (
    "public inherited sharing class SOQL_Account\n"
    "    extends SOQL implements SOQL.Selector {\n"
    "\n"
    "  public static SOQL_Account query() {\n"
    "    return new SOQL_Account();\n"
    "  }\n"
    "\n"
    "  private SOQL_Account() {\n"
    "    super(Account.SObjectType);\n"
    "    with(Account.Id, Account.Name);\n"
    "  }\n"
    "\n"
    "  public SOQL_Account byIndustry(String v) {\n"
    "    whereAre(\n"
    "      Filter.with(Account.Industry).equal(v));\n"
    "    return this;\n"
    "  }\n"
    "}"
)
add_text_box(slide, Inches(0.8), Inches(2.2), Inches(5.4), Inches(4.5),
             code_text, font_size=15, color=ACCENT_GREEN, font_name="Consolas")

# Right: What the agent learns
add_card_with_title(slide, Inches(6.9), Inches(1.6), Inches(5.9), Inches(5.4),
                    "What the Agent Learns", [
                        "Extend SOQL, implement SOQL.Selector",
                        "Static query() factory method",
                        "Default fields in constructor",
                        "Chainable filter methods",
                        "Return this for fluency",
                        "One example = every selector",
                        "the agent writes follows this pattern",
                    ], accent=ACCENT_ORANGE, font_size=17)


# ─── SLIDE 10: The Autonomous Pipeline ───────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.5), ACCENT_PURPLE)
add_text_box(slide, Inches(0.5), Inches(0.4), Inches(11), Inches(0.8),
             "   The Autonomous Pipeline", font_size=40, color=WHITE, bold=True)
add_text_box(slide, Inches(0.5), Inches(1.1), Inches(11), Inches(0.5),
             "   /story-to-feature \u2014 six phases, zero intervention",
             font_size=20, color=MID_GRAY)

phases = [
    (0, "Grill", "Stress-test the story\nfor gaps & edge cases", RGBColor(0x00, 0xBC, 0xD4)),
    (1, "Architect", "Design data model\nand plan all artifacts", ACCENT_BLUE),
    (2, "Implement", "Write Apex, LWC,\nmetadata", ACCENT_GREEN),
    (3, "Test", "Deploy, run, self-heal\n(max 3 iterations)", ACCENT_ORANGE),
    (4, "Validate", "Browser automation\nvia Playwright", ACCENT_PURPLE),
    (5, "Ship", "Review, commit,\npush, open PR", ACCENT_RED),
]

for i, (num, title, desc, color) in enumerate(phases):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + Inches(4.1) * col
    top = Inches(1.8) + Inches(2.8) * row
    add_phase_card(slide, left, top, Inches(3.8), num, title, desc, color)


# ─── SLIDE 11: Skills Architecture ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.5), ACCENT_BLUE)
add_text_box(slide, Inches(0.5), Inches(0.4), Inches(11), Inches(0.8),
             "   Skills \u2014 Reusable Agent Capabilities", font_size=40, color=WHITE, bold=True)
add_text_box(slide, Inches(0.5), Inches(1.2), Inches(11), Inches(0.5),
             "   Each skill is a prompt template that Claude Code can invoke",
             font_size=20, color=MID_GRAY)

# Orchestrator banner
orch_top = Inches(1.9)
add_card(slide, Inches(0.5), orch_top, Inches(12.3), Inches(1.1))
add_text_box(slide, Inches(0.8), orch_top + Inches(0.15), Inches(5.0), Inches(0.5),
             "/story-to-feature", font_size=26, color=ACCENT_ORANGE, bold=True,
             font_name="Consolas")
add_text_box(slide, Inches(0.8), orch_top + Inches(0.65), Inches(11.5), Inches(0.45),
             "Orchestrates all skills below in sequence \u2014 one command, full pipeline",
             font_size=16, color=LIGHT_GRAY)

# Six sub-skills in 2x3 grid
sub_skills = [
    ("/grill-me", "Self-grills the story\nfor gaps & edge cases", RGBColor(0x00, 0xBC, 0xD4)),
    ("/architect", "Designs the solution \u2014\ndata model, layers, plan", ACCENT_BLUE),
    ("/implement", "Writes all code following\nCLAUDE.md conventions", ACCENT_GREEN),
    ("/test", "Writes + runs tests,\nself-heals on failure", ACCENT_ORANGE),
    ("/validate", "Browser testing via\nPlaywright automation", ACCENT_PURPLE),
    ("/review", "Code review against\nbest practices", ACCENT_RED),
]

grid_top = Inches(3.3)
for i, (name, desc, color) in enumerate(sub_skills):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + Inches(4.1) * col
    top = grid_top + Inches(2.0) * row
    add_card(slide, left, top, Inches(3.8), Inches(1.8))
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2),
                 Inches(3.2), Inches(0.4), name,
                 font_size=22, color=color, bold=True, font_name="Consolas")
    add_text_box(slide, left + Inches(0.3), top + Inches(0.8),
                 Inches(3.2), Inches(0.9), desc,
                 font_size=14, color=LIGHT_GRAY)


# ─── SLIDE 12: Self-Healing ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.5), ACCENT_ORANGE)
add_text_box(slide, Inches(0.5), Inches(0.4), Inches(11), Inches(0.8),
             "   Self-Healing \u2014 The Secret Sauce", font_size=40, color=WHITE, bold=True)

# Left side - the loop
add_card(slide, Inches(0.5), Inches(1.6), Inches(6.0), Inches(5.4))
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.4), Inches(0.5),
             "Deploy \u2192 Test \u2192 Fail \u2192 Fix \u2192 Repeat", font_size=22,
             color=ACCENT_ORANGE, bold=True)

loop_items = [
    "\u2716  Deploy fails: \"Method does not exist: mockId\"",
    "\u2714  Agent reads error, finds correct API: .identifier()",
    "",
    "\u2716  Apex test fails: AuraHandledException message quirk",
    "\u2714  Agent adjusts assertion to match SF behavior",
    "",
    "\u2716  Jest fails: lwc:else not rendering in test env",
    "\u2714  Agent refactors to explicit lwc:if with computed getters",
]
add_bullet_list(slide, Inches(0.8), Inches(2.6), Inches(5.4), Inches(4.0),
                loop_items, font_size=15, color=LIGHT_GRAY, bullet_color=ACCENT_ORANGE)

# Right side - stats
add_card(slide, Inches(6.9), Inches(1.6), Inches(5.9), Inches(5.4))
add_text_box(slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(0.5),
             "From Today's Run", font_size=22, color=ACCENT_GREEN, bold=True)

stats = [
    "Deploy attempts: 2 (1 failed, 1 passed)",
    "Apex test runs: 2 (1 failure fixed)",
    "Jest test runs: 2 (3 failures fixed)",
    "Total self-heal cycles: 3",
    "",
    "Final result:",
    "  All Apex tests passing",
    "  All 12 Jest tests passing",
    "  Browser validation successful",
    "  31 files, 1,784 lines of code",
]
add_bullet_list(slide, Inches(7.2), Inches(2.6), Inches(5.3), Inches(4.0),
                stats, font_size=15, color=LIGHT_GRAY, bullet_color=ACCENT_GREEN)


# ─── SLIDE 13: Browser Validation ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.5), ACCENT_PURPLE)
add_text_box(slide, Inches(0.5), Inches(0.4), Inches(11), Inches(0.8),
             "   Browser Validation via Playwright", font_size=40, color=WHITE, bold=True)

add_card(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.4))

steps = [
    ("1. Login", "Gets org URL with session token via SF CLI", ACCENT_BLUE),
    ("2. Create Test Data", "Creates an Invoice record (Status: Sent, Amount: $1,000)", ACCENT_GREEN),
    ("3. Record Payment", "Fills form: $600, Credit Card, TXN-12345 \u2192 clicks Record Payment", ACCENT_ORANGE),
    ("4. Verify Update", "Confirms remaining balance updated to $400, payment in history table", ACCENT_PURPLE),
    ("5. Complete Payment", "Records $400 more \u2192 verifies \"fully paid\" banner + Status: Paid", ACCENT_RED),
    ("6. Screenshot Evidence", "Captures screenshots at each step as proof", RGBColor(0x00, 0xBC, 0xD4)),
]

for i, (step, desc, color) in enumerate(steps):
    top = Inches(1.9) + Inches(0.8) * i
    add_text_box(slide, Inches(1.0), top, Inches(3.0), Inches(0.6),
                 step, font_size=20, color=color, bold=True)
    add_text_box(slide, Inches(4.2), top, Inches(8.0), Inches(0.6),
                 desc, font_size=18, color=LIGHT_GRAY)


# ─── SLIDE 14: Let's Check In ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
             "Let's Check on Our Developer",
             font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             "\U0001F449  Switch to terminal",
             font_size=28, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 15: Key Takeaways ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.5), ACCENT_GREEN)
add_text_box(slide, Inches(0.5), Inches(0.4), Inches(11), Inches(0.8),
             "   Key Takeaways", font_size=40, color=WHITE, bold=True)

takeaways = [
    ("Conventions > Instructions", "A well-structured CLAUDE.md with examples teaches better than verbose rules", ACCENT_BLUE),
    ("Libraries Are Force Multipliers", "Consistent APIs mean the agent learns one pattern and applies it everywhere", ACCENT_GREEN),
    ("Mocking Is Non-Negotiable", "Full mockability enables autonomous testing without org dependencies", ACCENT_ORANGE),
    ("Self-Healing Makes It Practical", "Expect failures \u2014 the agent's ability to read errors and fix is the real power", ACCENT_PURPLE),
    ("Skills Compose", "Small, focused skills chain into powerful pipelines", ACCENT_RED),
]

for i, (title, desc, color) in enumerate(takeaways):
    top = Inches(1.6) + Inches(1.1) * i
    add_card(slide, Inches(0.5), top, Inches(12.3), Inches(0.95))
    add_text_box(slide, Inches(0.9), top + Inches(0.1), Inches(4.5), Inches(0.4),
                 title, font_size=22, color=color, bold=True)
    add_text_box(slide, Inches(0.9), top + Inches(0.5), Inches(11.5), Inches(0.4),
                 desc, font_size=17, color=LIGHT_GRAY)


# ─── SLIDE 16: Q&A ───────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1),
             "Questions?",
             font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(4.5), Inches(3.5), Inches(4), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
             "github.com/vernemi2/claude-code-user-group",
             font_size=22, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER,
             font_name="Consolas")

add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
             "Michal Verner",
             font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ─── Save ─────────────────────────────────────────────────────────────
output = "/Users/michal/Documents/personal/salesforce/repos/claude-code-user-group/docs/presentation.pptx"
prs.save(output)
print(f"Saved to {output}")
print(f"Total slides: {len(prs.slides)}")
